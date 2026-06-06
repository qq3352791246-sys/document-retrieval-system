import os
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import chardet


class DocumentReader:
    """文档读取器 - 支持Word、Excel、PPT"""
    
    @staticmethod
    def read_txt(filepath):
        """读取TXT文件"""
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        except:
            with open(filepath, 'r', encoding='gbk') as f:
                return f.read()
    
    @staticmethod
    def read_docx(filepath):
        """读取Word文档"""
        doc = Document(filepath)
        text = '\n'.join([para.text for para in doc.paragraphs])
        return text
    
    @staticmethod
    def read_excel(filepath):
        """读取Excel文件"""
        workbook = load_workbook(filepath)
        text = []
        for sheet in workbook.sheetnames:
            ws = workbook[sheet]
            text.append(f"=== Sheet: {sheet} ===\n")
            for row in ws.iter_rows(values_only=True):
                text.append(' '.join([str(cell) if cell else '' for cell in row]))
        return '\n'.join(text)
    
    @staticmethod
    def read_pptx(filepath):
        """读取PPT文件"""
        prs = Presentation(filepath)
        text = []
        for i, slide in enumerate(prs.slides):
            text.append(f"=== Slide {i+1} ===\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    
    @classmethod
    def read_document(cls, filepath):
        """根据文件类型读取文档"""
        _, ext = os.path.splitext(filepath)
        ext = ext.lower()
        
        if ext == '.txt':
            return cls.read_txt(filepath)
        elif ext == '.docx':
            return cls.read_docx(filepath)
        elif ext in ['.xlsx', '.xls']:
            return cls.read_excel(filepath)
        elif ext == '.pptx':
            return cls.read_pptx(filepath)
        else:
            raise ValueError(f"不支持的文件格式: {ext}")
    
    @classmethod
    def read_corpus(cls, corpus_path):
        """读取语料库中的所有文档"""
        documents = {}
        if not os.path.exists(corpus_path):
            os.makedirs(corpus_path)
            return documents
        
        for filename in os.listdir(corpus_path):
            if filename.startswith('.'):
                continue
            filepath = os.path.join(corpus_path, filename)
            try:
                content = cls.read_document(filepath)
                documents[filename] = content
            except Exception as e:
                print(f"读取文件 {filename} 失败: {e}")
        
        return documents